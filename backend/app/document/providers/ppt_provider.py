from __future__ import annotations

import os
from typing import Any

from app.domain.capability_types import CapabilityType
from app.domain.document_types import DocumentType
from app.document.providers.base import BaseDocumentProvider, ProviderResult

try:
    from pptx import Presentation  # pyright: ignore[reportMissingImports]
except Exception:
    Presentation = None


class PptProvider(BaseDocumentProvider):

    document_type = DocumentType.PPT

    def supported_capabilities(self) -> set[CapabilityType]:
        if not Presentation:
            return set()
        return {
            CapabilityType.READ,
            CapabilityType.EXTRACT,
            CapabilityType.LOCATE,
            CapabilityType.FILL,
            CapabilityType.VALIDATE,
            CapabilityType.WRITE,
            CapabilityType.SCAN_TEMPLATE,
        }

    def read(self, file_path: str, **kwargs) -> ProviderResult:
        if not Presentation:
            return ProviderResult(success=False, message="python-pptx is not available")
        try:
            pres = Presentation(file_path)
            preview: list[dict[str, Any]] = []
            for slide_idx, slide in enumerate(pres.slides):
                if slide_idx >= 5:
                    break
                texts = self._collect_slide_texts(slide)
                preview.append(
                    {
                        "slide_index": slide_idx,
                        "text_preview": texts[:10],
                    }
                )
            return ProviderResult(
                success=True,
                message="PPT document read successfully",
                data={
                    "slide_count": len(pres.slides),
                    "preview": preview,
                },
                raw={"provider": "ppt"},
            )
        except Exception as e:
            return ProviderResult(success=False, message=f"Failed to read ppt document: {e}")

    def extract(self, file_path: str, **kwargs) -> ProviderResult:
        if not Presentation:
            return ProviderResult(success=False, message="python-pptx is not available")
        try:
            pres = Presentation(file_path)
            slides: list[dict[str, Any]] = []
            for slide_idx, slide in enumerate(pres.slides):
                texts = self._collect_slide_texts(slide)
                tables = self._collect_slide_tables(slide)
                slides.append(
                    {
                        "slide_index": slide_idx,
                        "texts": texts,
                        "tables": tables,
                    }
                )
            return ProviderResult(
                success=True,
                message="PPT structured data extracted successfully",
                data={
                    "slide_count": len(slides),
                    "slides": slides,
                },
                raw={"provider": "ppt"},
            )
        except Exception as e:
            return ProviderResult(success=False, message=f"Failed to extract ppt data: {e}")

    def locate(self, file_path: str, **kwargs) -> ProviderResult:
        keyword = kwargs.get("keyword")
        if not keyword:
            return ProviderResult(success=False, message="'keyword' is required for ppt locate")
        if not Presentation:
            return ProviderResult(success=False, message="python-pptx is not available")

        try:
            pres = Presentation(file_path)
            case_sensitive = bool(kwargs.get("case_sensitive", False))
            needle = str(keyword) if case_sensitive else str(keyword).lower()

            matches: list[dict[str, Any]] = []
            for slide_idx, slide in enumerate(pres.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if not hasattr(shape, "text"):
                        continue
                    text = shape.text or ""
                    compare = text if case_sensitive else text.lower()
                    if needle in compare:
                        matches.append(
                            {
                                "slide_index": slide_idx,
                                "shape_index": shape_idx,
                                "text": text,
                            }
                        )

            return ProviderResult(
                success=True,
                message="PPT locate completed",
                data={
                    "keyword": keyword,
                    "matches": matches,
                },
                raw={"provider": "ppt"},
            )
        except Exception as e:
            return ProviderResult(success=False, message=f"Failed to locate in ppt: {e}")

    def fill(self, file_path: str, **kwargs) -> ProviderResult:
        if not Presentation:
            return ProviderResult(success=False, message="python-pptx is not available")
        try:
            pres = Presentation(file_path)
            field_values: dict[str, Any] = kwargs.get("field_values", {})
            replace_count = 0

            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        original = shape.text or ""
                        new_text, count = self._replace_placeholders_in_text(original, field_values)
                        replace_count += count
                        if new_text != original:
                            shape.text = new_text

                    if getattr(shape, "has_table", False):
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                original = cell.text or ""
                                new_text, count = self._replace_placeholders_in_text(original, field_values)
                                replace_count += count
                                if new_text != original:
                                    cell.text = new_text

            output_path = kwargs.get("output_path") or self._default_output_path(file_path)
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            pres.save(output_path)

            return ProviderResult(
                success=True,
                message="PPT document filled successfully",
                data={
                    "field_values_count": len(field_values),
                    "replace_count": replace_count,
                },
                output_path=output_path,
                raw={"provider": "ppt"},
            )
        except Exception as e:
            return ProviderResult(success=False, message=f"Failed to fill ppt document: {e}")

    def validate(self, file_path: str, **kwargs) -> ProviderResult:
        if not Presentation:
            return ProviderResult(success=False, message="python-pptx is not available")
        try:
            pres = Presentation(file_path)
            min_slides = int(kwargs.get("min_slides", 1))
            required_keywords: list[str] = kwargs.get("required_keywords", [])

            all_text = "\n".join(
                text
                for slide in pres.slides
                for text in self._collect_slide_texts(slide)
            )
            missing_keywords = [key for key in required_keywords if key not in all_text]

            valid = len(pres.slides) >= min_slides and not missing_keywords
            return ProviderResult(
                success=True,
                message="PPT validation completed",
                data={
                    "valid": valid,
                    "slide_count": len(pres.slides),
                    "min_slides": min_slides,
                    "required_keywords": required_keywords,
                    "missing_keywords": missing_keywords,
                },
                raw={"provider": "ppt"},
            )
        except Exception as e:
            return ProviderResult(success=False, message=f"Failed to validate ppt document: {e}")

    def write(self, file_path: str, **kwargs) -> ProviderResult:
        if not Presentation:
            return ProviderResult(success=False, message="python-pptx is not available")
        try:
            if kwargs.get("field_values"):
                return self.fill(file_path=file_path, **kwargs)

            pres = Presentation(file_path)
            output_path = kwargs.get("output_path") or self._default_output_path(file_path, suffix="_written")
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            pres.save(output_path)

            return ProviderResult(
                success=True,
                message="PPT document written successfully",
                data={"slide_count": len(pres.slides)},
                output_path=output_path,
                raw={"provider": "ppt"},
            )
        except Exception as e:
            return ProviderResult(success=False, message=f"Failed to write ppt document: {e}")

    def scan_template(self, file_path: str, **kwargs) -> ProviderResult:
        if not Presentation:
            return ProviderResult(success=False, message="python-pptx is not available")
        try:
            pres = Presentation(file_path)
            fields: set[str] = set()
            occurrences: list[dict[str, Any]] = []

            for slide_idx, slide in enumerate(pres.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if not hasattr(shape, "text"):
                        continue
                    text = shape.text or ""
                    shape_fields, shape_occurrences = self._scan_placeholders(text)
                    fields.update(shape_fields)
                    for item in shape_occurrences:
                        occurrences.append(
                            {
                                "slide_index": slide_idx,
                                "shape_index": shape_idx,
                                **item,
                            }
                        )

            return ProviderResult(
                success=True,
                message="PPT template scanned successfully",
                data={
                    "fields": sorted(fields),
                    "field_count": len(fields),
                    "occurrences": occurrences,
                },
                raw={"provider": "ppt"},
            )
        except Exception as e:
            return ProviderResult(success=False, message=f"Failed to scan ppt template: {e}")

    @staticmethod
    def _collect_slide_texts(slide) -> list[str]:
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            texts.append(cell.text)
        return texts

    @staticmethod
    def _collect_slide_tables(slide) -> list[dict[str, Any]]:
        tables: list[dict[str, Any]] = []
        for shape_idx, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_table", False):
                continue
            rows: list[list[str]] = []
            for row in shape.table.rows:
                rows.append([cell.text for cell in row.cells])
            tables.append({"shape_index": shape_idx, "rows": rows})
        return tables
